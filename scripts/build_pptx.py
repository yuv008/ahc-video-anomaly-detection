"""Build the 2-slide hackathon submission deck as a native .pptx.

Mirrors docs/slides.html content and palette exactly (dark ground, amber accent, teal=win,
red=bad) so the PPT and the published HTML artifact tell the identical story - this script
exists only because the submission rules require an actual PPT file, not a web page.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Palette - matches docs/slides.html :root dark theme exactly
GROUND = RGBColor(0x0E, 0x14, 0x1C)
SURFACE = RGBColor(0x16, 0x1F, 0x2B)
SURFACE2 = RGBColor(0x1D, 0x28, 0x36)
LINE = RGBColor(0x2A, 0x38, 0x49)
INK = RGBColor(0xE6, 0xED, 0xF6)
MUTED = RGBColor(0x84, 0x94, 0xA8)
AMBER = RGBColor(0xFF, 0x8A, 0x3D)
TEAL = RGBColor(0x45, 0xC2, 0xA0)
RED = RGBColor(0xFF, 0x5C, 0x5C)

SANS = "Calibri"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = GROUND
    return s


def box(slide, x, y, w, h, fill=None, line_color=LINE, line_w=0.75, top_accent=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill if fill else SURFACE2
    shp.line.color.rgb = line_color
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if top_accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.035))
        bar.fill.solid(); bar.fill.fore_color.rgb = top_accent
        bar.line.fill.background()
        bar.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, runs, size=11, font=SANS, color=INK, bold=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15, wrap=True):
    """runs: str, or list of (text, color, bold, size) tuples for one paragraph,
    or list of lists for multiple paragraphs."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0

    if isinstance(runs, str):
        runs = [[(runs, color, bold, size)]]
    elif runs and isinstance(runs[0], tuple):
        runs = [runs]

    for i, para_runs in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for rtext, rcolor, rbold, rsize in para_runs:
            r = p.add_run()
            r.text = rtext
            r.font.size = Pt(rsize)
            r.font.name = font
            r.font.color.rgb = rcolor
            r.font.bold = rbold
    return tb


def eyebrow(slide, s):
    text(slide, 0.5, 0.28, 10, 0.3, [(s.upper(), AMBER, False, 11)], font=MONO)


def footer(slide, left, right):
    ln = slide.shapes.add_connector(1, Inches(0.5), Inches(7.12), Inches(12.83), Inches(7.12))
    ln.line.color.rgb = LINE; ln.line.width = Pt(0.75)
    text(slide, 0.5, 7.18, 8, 0.3, [(left, MUTED, False, 9)], font=MONO)
    text(slide, 8.5, 7.18, 4.33, 0.3, [(right, MUTED, False, 9)], font=MONO, align=PP_ALIGN.RIGHT)


def table(slide, x, y, w, h, headers, rows, col_widths, font_size=10.5):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gshape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gshape.table
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = Inches(cw)

    # strip default banding style
    tbl_pr = tbl._tbl.find(qn('a:tblPr'))
    if tbl_pr is not None:
        tbl_pr.set('firstRow', '0')
        tbl_pr.set('bandRow', '0')

    def style_cell(cell, runs, bg=SURFACE2, align=PP_ALIGN.LEFT, size=font_size):
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        if isinstance(runs, str):
            runs = [(runs, INK, False)]
        for rtext, rcolor, rbold in runs:
            r = p.add_run()
            r.text = rtext
            r.font.size = Pt(size)
            r.font.name = SANS
            r.font.color.rgb = rcolor
            r.font.bold = rbold

    for j, htext in enumerate(headers):
        style_cell(tbl.cell(0, j), [(htext.upper(), MUTED, True)], bg=SURFACE2, size=font_size - 1.5)
    for i, row in enumerate(rows, start=1):
        for j, cellval in enumerate(row):
            if isinstance(cellval, tuple):
                runs, align = cellval
            else:
                runs, align = cellval, PP_ALIGN.LEFT
            style_cell(tbl.cell(i, j), runs, bg=SURFACE, align=align)
    return tbl


# ============================================================= SLIDE 1 =============================================================
s1 = add_slide()
eyebrow(s1, "AHC Visual Intelligence Hackathon · Slide 1 of 2 — What we built")
text(s1, 0.5, 0.56, 12.3, 0.98,
     [[("A small VLM decides windows; a second model's confidence filters false alarms; "
        "a cheap aggregator turns both into events", INK, True, 19)]],
     font=SANS, line_spacing=1.08)
text(s1, 0.5, 1.62, 12.3, 0.3,
     [("Nothing large runs at inference. Every number below is measured on the actual Tesla T4, not assumed.",
       MUTED, False, 11.5)])

# pipeline chain - 4 stages
stage_data = [
    ("STAGE 0", "Window sampler", "4s windows, no overlap, 8 frames. 50% overlap measured strictly "
     "worse: same F1, worse boundary IoU, 2x compute.", "4s budget / window / feed", False),
    ("STAGE 2", "Qwen3-VL-8B - 4-bit LoRA", "LoRA r=32, vision tower frozen, fp16. Emits minimal JSON; "
     "scores read from token logprobs, not the string.", "~4.3s / window on T4", True),
    ("STAGE 2.5", "Cosmos-Embed1 confidence gate", "VETOES Qwen's \"anomalous\" calls it isn't itself "
     ">=95% confident about. A filter on Qwen, not a second classifier.", "shipped - +9% score", True),
    ("STAGE 3", "Temporal aggregator", "Hysteresis + merge-gap + midpoint boundary shrink turns window "
     "verdicts into a clean event list.", "CPU only, free to re-tune", False),
]
sx, sy, sw, sh, gap = 0.5, 2.05, 2.98, 1.42, 0.1
for i, (n, t, d, cost, hot) in enumerate(stage_data):
    x = sx + i * (sw + gap)
    box(s1, x, sy, sw, sh, top_accent=(AMBER if hot else LINE))
    text(s1, x + 0.12, sy + 0.08, sw - 0.24, 0.2, [(n, MUTED, False, 9)], font=MONO)
    text(s1, x + 0.12, sy + 0.27, sw - 0.24, 0.35, [(t, INK, True, 13)])
    text(s1, x + 0.12, sy + 0.62, sw - 0.24, 0.65, [(d, MUTED, False, 9.5)], line_spacing=1.15)
    text(s1, x + 0.12, sy + sh - 0.3, sw - 0.24, 0.25, [(cost, TEAL, False, 9.5)], font=MONO)

# bottom two panels
py, ph = 3.62, 3.35
p1x, p1w = 0.5, 7.15
p2x, p2w = 7.78, 5.05

box(s1, p1x, py, p1w, ph)
text(s1, p1x + 0.18, py + 0.14, p1w - 0.36, 0.3,
     [("THE TRAP THE DATA SETS — WHY TIMESTAMPS COME FROM THE AGGREGATOR", MUTED, True, 10.5)], font=MONO)

# mini bar viz: train clip (one solid amber bar) vs test L2/L3 (sparse amber segments)
bar_y0 = py + 0.48
bar_x, bar_w = p1x + 1.1, p1w - 1.35
text(s1, p1x + 0.18, bar_y0, 0.9, 0.25, [("train clip", MUTED, False, 9)], font=MONO, align=PP_ALIGN.RIGHT)
b1 = box(s1, bar_x, bar_y0, bar_w, 0.22, fill=SURFACE)
seg = box(s1, bar_x, bar_y0, bar_w, 0.22, fill=AMBER, line_color=AMBER)

text(s1, p1x + 0.18, bar_y0 + 0.35, 0.9, 0.25, [("test L2/L3", MUTED, False, 9)], font=MONO, align=PP_ALIGN.RIGHT)
b2 = box(s1, bar_x, bar_y0 + 0.35, bar_w, 0.22, fill=SURFACE)
for (frac_l, frac_w) in [(0.08, 0.05), (0.27, 0.08), (0.62, 0.04), (0.80, 0.06)]:
    box(s1, bar_x + bar_w * frac_l, bar_y0 + 0.35, bar_w * frac_w, 0.22, fill=AMBER, line_color=AMBER)

text(s1, p1x + 0.18, bar_y0 + 0.68, p1w - 0.36, 2.0, [
    [("Training clips are pre-trimmed: the event fills ", MUTED, False, 10.5),
     ("99.9%", INK, True, 10.5), (" of the median clip. Real Level-2/3 footage is ", MUTED, False, 10.5),
     ("7.7%", INK, True, 10.5), (" anomaly. Trained naively, a model learns \"a clip was shown to "
      "me => something is wrong\" and fires continuously.", MUTED, False, 10.5)],
    [("", MUTED, False, 4)],
    [("Fix: window labels come from ", MUTED, False, 10.5), ("temporal overlap", INK, True, 10.5),
     (" with the annotation, not the folder a clip sits in - recovered 1,604 genuine background "
      "windows from inside anomaly clips. The model never predicts timestamps at all (no sub-clip "
      "localisation signal exists in training data to learn that from).", MUTED, False, 10.5)],
], line_spacing=1.25)

box(s1, p2x, py, p2w, ph)
text(s1, p2x + 0.18, py + 0.14, p2w - 0.36, 0.3,
     [("MODEL CHOICES — REASONED, NOT ASSUMED", MUTED, True, 10.5)], font=MONO)

model_rows = [
    ([("Qwen2.5-VL-3B", INK, False)], [("rejected", RED, True)], [("slower than 7B: 36 layers vs 28, depth dominates at batch 1", MUTED, False)]),
    ([("bf16 precision", INK, False)], [("rejected", RED, True)], [("T4 emulates it: 2.28 vs fp16 20.98 TFLOP/s - 9x slower", MUTED, False)]),
    ([("720p uncapped frames", INK, False)], [("rejected", RED, True)], [("1,196 tok/frame -> 17.8s/window; 256 cap = 5x speedup", MUTED, False)]),
    ([("Qwen3.5 (used by #1 team)", INK, False)], [("rejected", RED, True)], [("no 4-bit trainer build; gap to our pick is 0.0009 - noise", MUTED, False)]),
    ([("Qwen3-VL-8B, 4-bit LoRA", INK, True)], [("SHIPPED", TEAL, True)], [("every top AI-City-2026 team runs a ~8B Qwen VLM", MUTED, False)]),
]
table(s1, p2x + 0.18, py + 0.5, p2w - 0.36, ph - 0.7,
      ["Tried", "Verdict", "Why"], model_rows, col_widths=[1.55, 0.95, 2.02], font_size=9.5)

footer(s1, "Qwen3-VL-8B · 4-bit · LoRA r=32 · fp16 · Cosmos-Embed1 gate · Tesla T4 16GB",
       "no hosted model at inference")

# ============================================================= SLIDE 2 =============================================================
s2 = add_slide()
eyebrow(s2, "Slide 2 of 2 — Experiments, failures, and what actually moved the score")
text(s2, 0.5, 0.56, 12.3, 0.98,
     [[('Score tripled on the real platform — but 2 of 3 "improvements" that looked good '
        'locally failed on upload', INK, True, 19)]],
     line_spacing=1.08)
text(s2, 0.5, 1.62, 12.3, 0.3,
     [("Consistent root cause: our local metric under-punishes false alarms. The platform doesn't.",
       MUTED, False, 11.5)])

# KPI row
kpis = [
    ("29.4", "v1 — first real submission, uniform thresholds", AMBER),
    ("39.3", "v2 — per-level policy (loose L1, strict L2/3)", AMBER),
    ("43.4", "v3 — + Cosmos confidence gate, SHIPPED", TEAL),
    ("0.669", "measured ceiling once class-naming is fixed", AMBER),
]
kx, ky, kw, kh, kgap = 0.5, 2.05, 2.98, 0.85, 0.1
for i, (v, k, accent) in enumerate(kpis):
    x = kx + i * (kw + kgap)
    kb = box(s2, x, ky, kw, kh)
    bar = slide_bar = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(ky), Inches(0.04), Inches(kh))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background(); bar.shadow.inherit = False
    text(s2, x + 0.15, ky + 0.08, kw - 0.3, 0.4, [(v, INK, True, 22)], font=MONO)
    text(s2, x + 0.15, ky + 0.5, kw - 0.3, 0.42, [(k, MUTED, False, 9)], line_spacing=1.2)

py2, ph2 = 3.15, 3.82
p1x2, p1w2 = 0.5, 6.55
p2x2, p2w2 = 7.23, 5.6

box(s2, p1x2, py2, p1w2, ph2)
text(s2, p1x2 + 0.18, py2 + 0.14, p1w2 - 0.36, 0.3,
     [("SAME SECOND MODEL, TWO ROLES — ONLY ONE SURVIVED THE PLATFORM", MUTED, True, 10)], font=MONO)
text(s2, p1x2 + 0.18, py2 + 0.55, p1w2 - 0.36, 1.5, [
    [("SHIPPED  ", TEAL, True, 10.5), ("Cosmos as a confidence gate: ", INK, True, 10.5),
     ("only trust Qwen's \"anomalous\" call if Cosmos independently clears 95% confidence on "
      "that exact window.", MUTED, False, 10.5)],
    [("L3 0.150 -> ", MUTED, False, 10.5), ("0.249", TEAL, True, 10.5),
     (", overall 0.367 -> ", MUTED, False, 10.5), ("0.400", TEAL, True, 10.5),
     (" — one messy blob split into a real IoU-0.82 match, 7 noisy fragments filtered to 2, "
      "one mistimed event removed clean.", MUTED, False, 10.5)],
], line_spacing=1.3)
text(s2, p1x2 + 0.18, py2 + 2.2, p1w2 - 0.36, 1.7, [
    [("FAILED LIVE  ", RED, True, 10.5), ("Cosmos as the Level-1 classifier: ", INK, True, 10.5),
     ("measured better locally (0.617 vs Qwen's 0.550, both components up) — ", MUTED, False, 10.5),
     ("uploaded anyway, score dropped.", RED, True, 10.5)],
    [("Cosmos fired on 22/24 videos and wrecked 3 of 4 true-normal ones. Third time a local-only "
      "gain failed live; now we check false-alarm count directly before trusting any local number.",
      MUTED, False, 10.5)],
], line_spacing=1.3)

box(s2, p2x2, py2, p2w2, ph2)
text(s2, p2x2 + 0.18, py2 + 0.14, p2w2 - 0.36, 0.3,
     [("THE DEEPEST FINDING: LOCALISATION IS SOLVED, NAMING THE CLASS IS THE CEILING", MUTED, True, 10)], font=MONO)

class_rows = [
    ([("traffic_congestion", INK, False)], [("100%", TEAL, True)], [("100%", TEAL, True)]),
    ([("traffic_accident", INK, False)], [("10%", INK, False)], [("0%", RED, True)]),
    ([("loitering", INK, False)], [("0%", RED, True)], [("0%", RED, True)]),
]
table(s2, p2x2 + 0.18, py2 + 0.5, p2w2 - 0.36, 1.15,
      ["True class", "Model fires", "Correct"], class_rows,
      col_widths=[2.1, 1.35, 1.35], font_size=10)

text(s2, p2x2 + 0.18, py2 + 1.85, p2w2 - 0.36, 1.9, [
    [("An unsupervised, zero-training novelty score ", INK, True, 10.5),
     ("(how different a window looks from the rest of its own video) reaches ", MUTED, False, 10.5),
     ("AUC 0.864", TEAL, True, 10.5), (" for accidents — where the fine-tuned model scores ", MUTED, False, 10.5),
     ("0/79", RED, True, 10.5), (" on the same windows.", MUTED, False, 10.5)],
    [("", MUTED, False, 4)],
    [("With the class assumed known, that localiser alone reaches ", MUTED, False, 10.5),
     ("8 matched events", INK, True, 10.5), (" (vs our shipped 1). Root cause: trained on ", MUTED, False, 10.5),
     ("560 of 14,157", RED, True, 10.5),
     (" labelled windows available (~4%) — the fix is more training data, not more tuning.", MUTED, False, 10.5)],
], line_spacing=1.3)

footer(s2, "windows labelled by temporal overlap · scores from token logprobs · best-of scoring confirmed",
       "next: retrain on the full dataset")

OUT = "docs/AHC_Presentation.pptx"
prs.save(OUT)
print(f"wrote {OUT}")
